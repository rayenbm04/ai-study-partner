"""These build small real files in memory (PDF via reportlab, DOCX/PPTX/XLSX
via their own write APIs) and run the actual extractor against them — not
mocks. If pdfplumber, python-docx, python-pptx, or openpyxl ever change how
they expose text, these fail for real instead of a mock quietly agreeing.

Anything that would need a real vision-model call (scanned PDF pages, PPTX
images, standalone images) uses FakeLLMProvider instead — still exercises the
real image-rendering/extraction code path, just not a real network call."""
import io

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from app.core.exceptions import ExtractionError, UnsupportedFileTypeError
from app.services.knowledge_base.extractors.docx import extract_docx
from app.services.knowledge_base.extractors.image import extract_image
from app.services.knowledge_base.extractors.pdf import _fix_pdf_text, extract_pdf
from app.services.knowledge_base.extractors.pptx import extract_pptx
from app.services.knowledge_base.extractors.registry import SUPPORTED_EXTENSIONS, extension_of, get_extractor
from app.services.knowledge_base.extractors.txt import extract_txt
from app.services.knowledge_base.extractors.xls import extract_xls
from app.services.knowledge_base.extractors.xlsx import extract_xlsx
from tests.unit.fakes import FakeLLMProvider


def _make_pdf(pages: list[str]) -> bytes:
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    for text in pages:
        c.drawString(72, 700, text)
        c.showPage()
    c.save()
    return buffer.getvalue()


def _make_blank_pdf_page() -> bytes:
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    c.showPage()  # no drawString call at all -> pdfplumber extracts no text
    c.save()
    return buffer.getvalue()


def _make_pdf_with_text_and_image(text: str) -> bytes:
    """A page with plenty of extracted text (well over the vision-fallback
    threshold) that ALSO has a large embedded image — the case
    _has_significant_image exists for, distinct from the low-text/scanned-page
    case _MIN_TEXT_CHARS_BEFORE_VISION_FALLBACK covers."""
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    c.drawString(72, 700, text)
    diagram = Image.new("RGB", (100, 100), color="blue")
    # 300x300pt on a 612x792pt letter page covers ~11.6% of the page area,
    # comfortably over _MIN_IMAGE_AREA_RATIO (3%).
    c.drawImage(ImageReader(diagram), 100, 300, width=300, height=300)
    c.showPage()
    c.save()
    return buffer.getvalue()


async def test_extract_pdf_returns_one_segment_per_page():
    content = _make_pdf(["First page of the course.", "Second page continues here."])
    segments = await extract_pdf(content, "notes.pdf")

    assert len(segments) == 2
    assert segments[0].page == 1
    assert "First page" in segments[0].text
    assert segments[1].page == 2
    assert "Second page" in segments[1].text


async def test_extract_pdf_skips_blank_pages_when_no_llm_provider():
    content = _make_pdf(["Only real content page.", ""])
    segments = await extract_pdf(content, "notes.pdf")
    assert len(segments) == 1


async def test_extract_pdf_raises_extraction_error_on_garbage_bytes():
    with pytest.raises(ExtractionError):
        await extract_pdf(b"this is not a pdf", "broken.pdf")


async def test_extract_pdf_falls_back_to_vision_for_low_text_page():
    content = _make_blank_pdf_page()
    llm = FakeLLMProvider(vision_response="Transcribed diagram content.")

    segments = await extract_pdf(content, "scanned.pdf", llm_provider=llm)

    assert len(segments) == 1
    assert segments[0].text == "Transcribed diagram content."
    assert len(llm.vision_calls) == 1
    assert llm.vision_calls[0]["mime_type"] == "image/png"


async def test_extract_pdf_does_not_call_vision_when_page_has_enough_text():
    content = _make_pdf(["This page already has plenty of real extracted text on it, well over fifty characters."])
    llm = FakeLLMProvider(vision_response="should not be used")

    await extract_pdf(content, "notes.pdf", llm_provider=llm)

    assert llm.vision_calls == []


async def test_extract_pdf_describes_significant_image_on_a_text_heavy_page():
    """A page with plenty of text AND a large embedded diagram should still
    get a vision call — just to describe the image, not replace the text —
    otherwise the diagram's content is silently dropped from the RAG index."""
    content = _make_pdf_with_text_and_image(
        "This page has plenty of real text content already, well over the fifty character threshold."
    )
    llm = FakeLLMProvider(vision_response="A free-body diagram showing forces on a block.")

    segments = await extract_pdf(content, "diagram.pdf", llm_provider=llm)

    assert len(llm.vision_calls) == 1
    assert len(segments) == 1
    assert "plenty of real text content" in segments[0].text
    assert "[Figure]: A free-body diagram" in segments[0].text


async def test_extract_pdf_skips_appending_when_vision_reports_no_meaningful_image():
    content = _make_pdf_with_text_and_image(
        "This page has plenty of real text content already, well over the fifty character threshold."
    )
    llm = FakeLLMProvider(vision_response="NONE")

    segments = await extract_pdf(content, "diagram.pdf", llm_provider=llm)

    assert len(llm.vision_calls) == 1
    assert "[Figure]" not in segments[0].text


async def test_extract_pdf_skips_vision_entirely_without_llm_provider_even_with_image():
    content = _make_pdf_with_text_and_image(
        "This page has plenty of real text content already, well over the fifty character threshold."
    )

    segments = await extract_pdf(content, "diagram.pdf", llm_provider=None)

    assert "[Figure]" not in segments[0].text


def test_fix_pdf_text_repairs_ligature_placeholders_and_mac_roman_mismatches():
    assert _fix_pdf_text("re(cid:12)ection") == "reflection"  # (cid:12) ligature placeholder -> "fl"
    assert _fix_pdf_text("(cid:99) stripped") == " stripped"  # unmapped cid codes are stripped, not left in
    assert _fix_pdf_text("cafØ") == "café"  # "Ø" MacRoman mismatch -> "é"


def _make_docx() -> bytes:
    document = DocxDocument()
    document.add_heading("Chapter 1: Vectors", level=1)
    document.add_paragraph("A vector has magnitude and direction.")
    document.add_heading("Chapter 2: Matrices", level=1)
    document.add_paragraph("A matrix is a rectangular array of numbers.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Term"
    table.rows[0].cells[1].text = "Definition"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


async def test_extract_docx_groups_paragraphs_under_headings():
    segments = await extract_docx(_make_docx(), "algebra.docx")
    section_titles = [s.section_title for s in segments]
    assert "Chapter 1: Vectors" in section_titles
    assert "Chapter 2: Matrices" in section_titles

    vectors_segment = next(s for s in segments if s.section_title == "Chapter 1: Vectors")
    assert "magnitude and direction" in vectors_segment.text


async def test_extract_docx_includes_table_content():
    segments = await extract_docx(_make_docx(), "algebra.docx")
    assert any("Term | Definition" in s.text for s in segments)


async def test_extract_docx_raises_extraction_error_on_garbage_bytes():
    with pytest.raises(ExtractionError):
        await extract_docx(b"not a docx file", "broken.docx")


def _make_pptx_with_image() -> bytes:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Newton's Laws"
    slide.placeholders[1].text = "The first law states that an object at rest stays at rest."

    tiny_image = Image.new("RGB", (4, 4), color="blue")
    image_buffer = io.BytesIO()
    tiny_image.save(image_buffer, format="PNG")
    image_buffer.seek(0)
    slide.shapes.add_picture(image_buffer, Inches(1), Inches(4), Inches(1), Inches(1))

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


async def test_extract_pptx_captures_title_and_body():
    segments = await extract_pptx(_make_pptx_with_image(), "physics.pptx")
    assert len(segments) == 1
    assert segments[0].section_title == "Newton's Laws"
    assert "object at rest" in segments[0].text
    assert segments[0].page == 1


async def test_extract_pptx_describes_embedded_images_when_llm_provided():
    llm = FakeLLMProvider(vision_response="A free-body diagram of a block on an incline.")

    segments = await extract_pptx(_make_pptx_with_image(), "physics.pptx", llm_provider=llm)

    assert len(llm.vision_calls) == 1
    assert "[Image]: A free-body diagram" in segments[0].text


async def test_extract_pptx_skips_image_description_without_llm_provider():
    segments = await extract_pptx(_make_pptx_with_image(), "physics.pptx")
    assert "[Image]" not in segments[0].text


async def test_extract_pptx_raises_extraction_error_on_garbage_bytes():
    with pytest.raises(ExtractionError):
        await extract_pptx(b"not a pptx file", "broken.pptx")


def _make_xlsx() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Formulas"
    sheet.append(["Name", "Formula"])
    sheet.append(["Area of circle", "pi * r^2"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


async def test_extract_xlsx_joins_cells_and_uses_sheet_name_as_section():
    segments = await extract_xlsx(_make_xlsx(), "formulas.xlsx")
    assert len(segments) == 1
    assert segments[0].section_title == "Formulas"
    assert "Area of circle | pi * r^2" in segments[0].text


async def test_extract_xlsx_raises_extraction_error_on_garbage_bytes():
    with pytest.raises(ExtractionError):
        await extract_xlsx(b"not an xlsx file", "broken.xlsx")


async def test_extract_xls_raises_extraction_error_on_garbage_bytes():
    with pytest.raises(ExtractionError):
        await extract_xls(b"not a real xls file", "broken.xls")


async def test_extract_txt_decodes_utf8():
    segments = await extract_txt("Résumé du cours de physique.".encode("utf-8"), "notes.txt")
    assert len(segments) == 1
    assert "Résumé" in segments[0].text


async def test_extract_txt_falls_back_to_latin1():
    content = "Café".encode("latin-1")
    segments = await extract_txt(content, "notes.txt")
    assert "Caf" in segments[0].text


async def test_extract_txt_empty_file_returns_no_segments():
    assert await extract_txt(b"   \n  ", "empty.txt") == []


async def test_extract_image_requires_llm_provider():
    with pytest.raises(ExtractionError):
        await extract_image(b"fake image bytes", "diagram.png", llm_provider=None)


async def test_extract_image_returns_vision_description():
    llm = FakeLLMProvider(vision_response="A hand-drawn circuit diagram with a resistor and battery.")

    segments = await extract_image(b"fake image bytes", "diagram.png", llm_provider=llm)

    assert len(segments) == 1
    assert "circuit diagram" in segments[0].text
    assert llm.vision_calls[0]["mime_type"] == "image/png"


async def test_extract_image_infers_mime_type_from_extension():
    llm = FakeLLMProvider(vision_response="A photo of handwritten notes.")
    await extract_image(b"fake jpeg bytes", "photo.JPG", llm_provider=llm)
    assert llm.vision_calls[0]["mime_type"] == "image/jpeg"


@pytest.mark.parametrize(
    "filename,expected_extension", [("Notes.PDF", ".pdf"), ("slides.pptx", ".pptx"), ("noextension", "")]
)
def test_extension_of(filename, expected_extension):
    assert extension_of(filename) == expected_extension


def test_get_extractor_returns_callable_for_supported_types():
    for extension in SUPPORTED_EXTENSIONS:
        extractor = get_extractor(f"file{extension}")
        assert callable(extractor)


def test_get_extractor_rejects_unsupported_type():
    with pytest.raises(UnsupportedFileTypeError):
        get_extractor("scan.gif")
