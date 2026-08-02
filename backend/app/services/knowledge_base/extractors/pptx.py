import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.exceptions import ExtractionError
from app.services.knowledge_base.extractors.base import ExtractedSegment
from app.services.llm.base import LLMProvider

_IMAGE_DESCRIPTION_PROMPT = (
    "Describe this image from a course slide in one or two sentences, focused on "
    "what a student would need to understand from it (a diagram's structure, a "
    "graph's trend, a formula shown in the image, etc.)."
)


async def extract_pptx(
    content: bytes, filename: str, llm_provider: LLMProvider | None = None
) -> list[ExtractedSegment]:
    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise ExtractionError(filename, str(exc)) from exc

    segments: list[ExtractedSegment] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_shape_id = title_shape.shape_id if title_shape is not None else None

        title: str | None = None
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title_shape_id is not None and shape.shape_id == title_shape_id:
                title = text
            else:
                body_parts.append(text)

        if llm_provider is not None:
            body_parts.extend(await _describe_images(slide, slide_number, filename, llm_provider))

        notes_text = _slide_notes(slide)
        if notes_text:
            body_parts.append(f"Speaker notes: {notes_text}")

        body = "\n".join(body_parts).strip()
        combined = f"{title}\n{body}".strip() if title else body
        if combined:
            segments.append(ExtractedSegment(text=combined, page=slide_number, section_title=title))

    return segments


async def _describe_images(slide, slide_number: int, filename: str, llm_provider: LLMProvider) -> list[str]:
    descriptions: list[str] = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            image = shape.image
            description = await llm_provider.complete_vision(
                image_bytes=image.blob,
                mime_type=image.content_type or "image/png",
                prompt=f"{_IMAGE_DESCRIPTION_PROMPT}\n\n(Slide {slide_number} of '{filename}'.)",
            )
            description = description.strip()
            if description:
                descriptions.append(f"[Image]: {description}")
        except Exception:
            # One bad image shouldn't fail extraction of the rest of the slide.
            continue
    return descriptions


def _slide_notes(slide) -> str:
    try:
        if slide.has_notes_slide:
            return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        pass
    return ""
